#!/usr/bin/env python3
"""Build interview-ppt assets from a paper document.

Outputs:
1) Structured outline markdown
2) Slide deck JSON (title + bullets + notes)
3) Google Workspace MCP batch_update payload template
4) Local draft PPTX (pandoc fallback)
5) Local styled PPTX (python-pptx, preferred)
"""

from __future__ import annotations

import argparse
import json
import re
import subprocess
import textwrap
from pathlib import Path
from typing import Dict, List


def run_command(args: List[str]) -> str:
    proc = subprocess.run(args, check=True, capture_output=True, text=True)
    return proc.stdout


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\u00a0", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_markdown(source: Path) -> str:
    suffix = source.suffix.lower()
    if suffix in {".docx", ".doc"}:
        md = run_command(["pandoc", str(source), "-t", "markdown"])
        return normalize_text(md)
    if suffix == ".pdf":
        txt = run_command(["pdftotext", "-layout", str(source), "-"])
        return normalize_text(txt)
    return normalize_text(source.read_text(encoding="utf-8"))


def first_title(markdown_text: str) -> str:
    lines = [ln.strip() for ln in markdown_text.splitlines() if ln.strip()]
    for idx, line in enumerate(lines):
        if set(line) <= {"=", "-"}:
            continue
        if idx + 1 < len(lines) and set(lines[idx + 1]) <= {"=", "-"}:
            return line
        if len(line) > 8:
            return line
    return "复试论文汇报"


def extract_abstract_parts(markdown_text: str) -> Dict[str, str]:
    # Try to isolate Chinese abstract segment.
    abstract_block = ""
    m = re.search(r"摘要[\s\S]{0,2500}?关键词", markdown_text)
    if m:
        abstract_block = m.group(0)
    else:
        # fallback to first 3000 chars
        abstract_block = markdown_text[:3000]

    def pick(label: str, stop_labels: List[str]) -> str:
        stop_pattern = "|".join(re.escape(x) for x in stop_labels)
        pattern = re.compile(rf"{re.escape(label)}[:：]\s*([\s\S]*?)(?=(?:{stop_pattern})[:：])")
        m2 = pattern.search(abstract_block)
        return normalize_text(m2.group(1)) if m2 else ""

    objective = pick("目的", ["方法", "结果", "结论"])
    methods = pick("方法", ["结果", "结论"])
    results = pick("结果", ["结论"])
    conclusion = ""
    m3 = re.search(r"结论[:：]\s*([\s\S]*?)(?:关键词|$)", abstract_block)
    if m3:
        conclusion = normalize_text(m3.group(1))

    kw = ""
    m4 = re.search(r"关键词[:：]\s*(.+)", abstract_block)
    if m4:
        kw = normalize_text(m4.group(1))

    return {
        "objective": objective,
        "methods": methods,
        "results": results,
        "conclusion": conclusion,
        "keywords": kw,
    }


def extract_core_numbers(markdown_text: str) -> Dict[str, str]:
    merged = markdown_text
    out: Dict[str, str] = {}

    # Sample size.
    m_n = re.search(r"(\d{2,4})例", merged)
    if m_n:
        out["sample_size"] = m_n.group(1)
    m_group = re.search(r"观察组[和与及 ]*对照组各(\d{1,4})例", merged)
    if m_group:
        out["group_size"] = m_group.group(1)

    # Satisfaction rate (prefer table row extraction to avoid "高于对照组16%" interference).
    sat_ctrl = ""
    sat_obs = ""
    for line in merged.splitlines():
        ln = line.strip()
        if "对照组" in ln and "%" in ln and "50" in ln:
            p = re.findall(r"([0-9]{1,3}%)", ln)
            if p:
                sat_ctrl = p[-1]
        if "观察组" in ln and "%" in ln and "50" in ln:
            p = re.findall(r"([0-9]{1,3}%)", ln)
            if p:
                sat_obs = p[-1]
    if sat_ctrl and sat_obs:
        out["satisfaction_control"] = sat_ctrl
        out["satisfaction_observation"] = sat_obs
    else:
        m_sat = re.search(r"观察组护理总满意度为([0-9]{1,3}%)[^。\n]*对照组([0-9]{1,3}%)", merged)
        if m_sat:
            out["satisfaction_observation"] = m_sat.group(1)
            out["satisfaction_control"] = m_sat.group(2)

    # VAS values from table 3.
    m_vas_ctrl = re.search(
        r"对照组\s+50\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)",
        merged,
    )
    m_vas_obs = re.search(
        r"观察组\s+50\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)",
        merged,
    )
    if m_vas_ctrl and m_vas_obs:
        out["vas_control"] = " / ".join(m_vas_ctrl.groups())
        out["vas_observation"] = " / ".join(m_vas_obs.groups())

    # Compliance values (table 6).
    m_compliance_ctrl = re.search(
        r"对照组\s+50\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)",
        merged,
    )
    m_compliance_obs = re.search(
        r"观察组\s+50\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)\s+([0-9.]+±[0-9.]+)",
        merged,
    )
    if m_compliance_ctrl and m_compliance_obs:
        out["compliance_control"] = " / ".join(m_compliance_ctrl.groups())
        out["compliance_observation"] = " / ".join(m_compliance_obs.groups())

    return out


def parse_mean(value: str) -> float:
    """Parse mean from strings like '5.24±2.66'."""
    m = re.search(r"([0-9]+(?:\.[0-9]+)?)", value)
    return float(m.group(1)) if m else 0.0


def parse_percent(value: str) -> float:
    m = re.search(r"([0-9]+(?:\.[0-9]+)?)\s*%", value)
    return float(m.group(1)) if m else 0.0


def extract_chart_data(numbers: Dict[str, str]) -> Dict[str, object]:
    chart_data: Dict[str, object] = {}

    vas_ctrl_raw = [x.strip() for x in numbers.get("vas_control", "").split("/") if x.strip()]
    vas_obs_raw = [x.strip() for x in numbers.get("vas_observation", "").split("/") if x.strip()]
    if len(vas_ctrl_raw) == 3 and len(vas_obs_raw) == 3:
        chart_data["vas"] = {
            "categories": ["术后当天", "术后第1天", "术后第1周"],
            "control": [parse_mean(x) for x in vas_ctrl_raw],
            "observation": [parse_mean(x) for x in vas_obs_raw],
        }

    sat_obs = parse_percent(numbers.get("satisfaction_observation", "100%"))
    sat_ctrl = parse_percent(numbers.get("satisfaction_control", "84%"))
    if sat_obs > 0 and sat_ctrl > 0:
        chart_data["satisfaction"] = {
            "categories": ["护理总满意度"],
            "control": [sat_ctrl],
            "observation": [sat_obs],
        }

    comp_ctrl_raw = [x.strip() for x in numbers.get("compliance_control", "").split("/") if x.strip()]
    comp_obs_raw = [x.strip() for x in numbers.get("compliance_observation", "").split("/") if x.strip()]
    if len(comp_ctrl_raw) == 5 and len(comp_obs_raw) == 5:
        chart_data["compliance"] = {
            "categories": ["情绪管理", "饮食控制", "戒烟戒酒", "及时用药", "功能锻炼"],
            "control": [parse_mean(x) for x in comp_ctrl_raw],
            "observation": [parse_mean(x) for x in comp_obs_raw],
        }

    return chart_data


def to_short(text: str, max_len: int = 90) -> str:
    clean = re.sub(r"\s+", " ", text).strip()
    if len(clean) <= max_len:
        return clean
    return clean[: max_len - 1] + "…"


def build_slides(
    title: str,
    abstract_parts: Dict[str, str],
    numbers: Dict[str, str],
    presenter: str,
    affiliation: str,
) -> List[Dict[str, object]]:
    n = numbers.get("sample_size", "100")
    group = numbers.get("group_size", "50")
    satisfaction_obs = numbers.get("satisfaction_observation", "100%")
    satisfaction_ctrl = numbers.get("satisfaction_control", "84%")
    vas_ctrl = numbers.get("vas_control", "5.24±2.66 / 4.77±2.45 / 3.12±1.03")
    vas_obs = numbers.get("vas_observation", "3.15±2.02 / 3.23±1.33 / 1.16±0.66")

    objective_short = to_short(abstract_parts.get("objective", "评估时效性激励护理在Rezum术后患者中的效果。"))
    method_short = to_short(abstract_parts.get("methods", "随机对照设计，对比常规护理与时效性激励护理干预。"))
    result_short = to_short(abstract_parts.get("results", "干预组在症状、疼痛、负性情绪、满意度与依从性方面整体优于对照组。"))
    conclusion_short = to_short(abstract_parts.get("conclusion", "时效性激励护理具有临床推广价值。"))
    keywords = abstract_parts.get("keywords", "时效性激励理论；前列腺增生；蒸汽消融术；护理干预")

    slides: List[Dict[str, object]] = [
        {
            "title": title,
            "bullets": [
                "投稿论文汇报（复试版）",
                f"汇报人：{presenter}",
                f"单位：{affiliation}",
            ],
            "notes": [
                "开场30秒说明研究背景与临床价值。",
                "一句话点明：这是可落地的术后护理改进方案。",
            ],
        },
        {
            "title": "研究问题与临床痛点",
            "bullets": [
                "BPH患者术后常见问题：疼痛、焦虑、控尿困难、依从性不足。",
                "Rezum疗效确切，但术后护理质量决定恢复速度与体验。",
                "传统护理难以覆盖关键时点的行为激励与心理支持。",
            ],
            "notes": ["强调“治疗 + 护理”共同决定结局。"],
        },
        {
            "title": "研究目的与核心假设",
            "bullets": [
                f"研究目的：{objective_short}",
                "核心假设：在关键时点给予精准激励，可改善症状与心理状态。",
                f"关键词：{keywords}",
            ],
            "notes": ["把“时效性”解释为干预时机匹配，而非单纯增加护理工作量。"],
        },
        {
            "title": "研究设计与样本",
            "bullets": [
                f"研究类型：随机对照研究（单中心），总样本 n={n}。",
                f"分组：观察组 {group} 例；对照组 {group} 例。",
                "观察指标：IPSS、控尿能力、VAS、SAS/SDS、满意度、依从性。",
            ],
            "notes": ["提及伦理批准号与知情同意，体现规范性。"],
        },
        {
            "title": "干预方案：时效性激励护理",
            "bullets": [
                "对照组：常规围术期护理。",
                "观察组：常规护理 + 需要激励 + 利益激励 + 情感激励 + 榜样激励。",
                "关键节点覆盖：术前-术后早期-居家恢复期。",
            ],
            "notes": ["突出“评估-反馈-激励”闭环，而非单次宣教。"],
        },
        {
            "title": "主要结果总览",
            "bullets": [
                f"结果摘要：{result_short}",
                "术后1个月：IPSS、SAS、SDS差异显著（文中报道 P<0.001）。",
                "术后1日/1周：VAS显著下降（文中报道 P<0.001）。",
            ],
            "notes": ["这一页只给方向性结论，下一页给关键数值。"],
        },
        {
            "title": "关键数值：疼痛与满意度",
            "bullets": [
                f"VAS（对照组，术后当天/第1天/第1周）：{vas_ctrl}",
                f"VAS（观察组，术后当天/第1天/第1周）：{vas_obs}",
                f"护理总满意度：观察组 {satisfaction_obs} vs 对照组 {satisfaction_ctrl}。",
            ],
            "notes": [
                "若最终投稿表格有更新，口头以最终表格为准。",
                "建议在正式版PPT插入柱状图提升可读性。",
            ],
        },
        {
            "title": "依从性与行为改变",
            "bullets": [
                "观察组在情绪管理、饮食控制、戒烟戒酒、及时用药、功能锻炼方面均优于对照组。",
                "提示干预不仅改善短期症状，也改善长期行为执行。",
                "临床意义：可复制到其他术后康复路径管理场景。",
            ],
            "notes": ["这页对应复试老师常问的“能否推广”。"],
        },
        {
            "title": "机制解释与临床价值",
            "bullets": [
                "机制层面：关键时机激励提升患者自我效能与康复参与度。",
                "护理价值：改善症状 + 缓解负性情绪 + 提升满意度与依从性。",
                "管理价值：形成标准化、可执行、可评估的护理闭环。",
            ],
            "notes": ["用“机制-结果-价值”三段式回答研究意义。"],
        },
        {
            "title": "研究局限与改进方向",
            "bullets": [
                "局限：单中心、样本量中等、随访时间较短。",
                "后续：多中心扩大样本，延长随访，强化量表与客观指标联动。",
                "计划：在硕士阶段延展到肿瘤护理与长期症状管理场景。",
            ],
            "notes": ["主动讲局限会加分，体现科研审慎性。"],
        },
        {
            "title": "复试可能提问与应答",
            "bullets": [
                "Q1：为什么选择时效性激励理论？",
                "A1：它直接面向“关键时点行为改变”，契合术后康复需求。",
                "Q2：如何控制偏倚？",
                "A2：随机分组、统一评估时点、量表与客观指标结合。",
                "Q3：如何与未来导师方向衔接？",
                "A3：将该干预框架迁移到肿瘤患者症状管理与心理支持。",
            ],
            "notes": ["最后30秒请老师提问，过渡自然。"],
        },
        {
            "title": "结论与致谢",
            "bullets": [
                f"结论：{conclusion_short}",
                "核心贡献：提出可执行的术后激励护理路径。",
                "感谢各位老师指正。",
            ],
            "notes": ["收束：贡献一句话 + 谢谢。"],
        },
    ]
    return slides


def build_workspace_requests(slides: List[Dict[str, object]]) -> List[Dict[str, object]]:
    requests: List[Dict[str, object]] = []
    for idx, slide in enumerate(slides, start=1):
        slide_id = f"s{idx:02d}"
        title_id = f"s{idx:02d}_title"
        body_id = f"s{idx:02d}_body"
        title_text = str(slide["title"])
        body_text = "\n".join(str(x) for x in slide["bullets"])

        requests.append(
            {
                "createSlide": {
                    "objectId": slide_id,
                    "insertionIndex": idx - 1,
                    "slideLayoutReference": {"predefinedLayout": "BLANK"},
                }
            }
        )
        requests.append(
            {
                "createShape": {
                    "objectId": title_id,
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": slide_id,
                        "size": {
                            "height": {"magnitude": 56, "unit": "PT"},
                            "width": {"magnitude": 640, "unit": "PT"},
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 40,
                            "translateY": 22,
                            "unit": "PT",
                        },
                    },
                }
            }
        )
        requests.append({"insertText": {"objectId": title_id, "insertionIndex": 0, "text": title_text}})
        requests.append(
            {
                "updateTextStyle": {
                    "objectId": title_id,
                    "textRange": {"type": "ALL"},
                    "style": {
                        "fontFamily": "Arial",
                        "fontSize": {"magnitude": 30, "unit": "PT"},
                        "bold": True,
                    },
                    "fields": "fontFamily,fontSize,bold",
                }
            }
        )
        requests.append(
            {
                "createShape": {
                    "objectId": body_id,
                    "shapeType": "TEXT_BOX",
                    "elementProperties": {
                        "pageObjectId": slide_id,
                        "size": {
                            "height": {"magnitude": 400, "unit": "PT"},
                            "width": {"magnitude": 640, "unit": "PT"},
                        },
                        "transform": {
                            "scaleX": 1,
                            "scaleY": 1,
                            "translateX": 40,
                            "translateY": 92,
                            "unit": "PT",
                        },
                    },
                }
            }
        )
        requests.append({"insertText": {"objectId": body_id, "insertionIndex": 0, "text": body_text}})
        requests.append(
            {
                "updateTextStyle": {
                    "objectId": body_id,
                    "textRange": {"type": "ALL"},
                    "style": {
                        "fontFamily": "Arial",
                        "fontSize": {"magnitude": 18, "unit": "PT"},
                    },
                    "fields": "fontFamily,fontSize",
                }
            }
        )
        requests.append(
            {
                "createParagraphBullets": {
                    "objectId": body_id,
                    "textRange": {"type": "ALL"},
                    "bulletPreset": "BULLET_DISC_CIRCLE_SQUARE",
                }
            }
        )
    return requests


def render_outline_markdown(
    title: str,
    presenter: str,
    affiliation: str,
    abstract_parts: Dict[str, str],
    slides: List[Dict[str, object]],
) -> str:
    lines: List[str] = []
    lines.append(f"# 论文投稿汇报PPT提纲\n")
    lines.append(f"- 论文标题：{title}")
    lines.append(f"- 汇报人：{presenter}")
    lines.append(f"- 单位：{affiliation}")
    lines.append("")
    lines.append("## 摘要提炼")
    lines.append(f"- 目的：{abstract_parts.get('objective', '')}")
    lines.append(f"- 方法：{abstract_parts.get('methods', '')}")
    lines.append(f"- 结果：{abstract_parts.get('results', '')}")
    lines.append(f"- 结论：{abstract_parts.get('conclusion', '')}")
    lines.append("")
    lines.append("## 幻灯片目录（12页）")
    for i, slide in enumerate(slides, start=1):
        lines.append(f"- 第{i}页：{slide['title']}")
    lines.append("")
    lines.append("## 使用提示")
    lines.append("- 建议汇报总时长 5-8 分钟，重点页为“结果总览”和“关键数值”。")
    lines.append("- 表格数值如后续稿件更新，请以最终投稿稿为准。")
    return "\n".join(lines).strip() + "\n"


def render_pandoc_slide_markdown(slides: List[Dict[str, object]]) -> str:
    parts: List[str] = []
    for slide in slides:
        parts.append(f"# {slide['title']}\n")
        for bullet in slide["bullets"]:
            parts.append(f"- {bullet}")
        notes = slide.get("notes", [])
        if notes:
            parts.append("")
            parts.append("::: notes")
            for note in notes:
                parts.append(f"- {note}")
            parts.append(":::")
        parts.append("")
    return "\n".join(parts).strip() + "\n"


def render_styled_pptx(
    slides: List[Dict[str, object]],
    out_path: Path,
    presenter: str,
    affiliation: str,
    chart_data: Dict[str, object] | None = None,
) -> None:
    """Render an academic-style PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # Palette: restrained academic style.
    C_BG = RGBColor(255, 255, 255)
    C_NAVY = RGBColor(21, 45, 80)
    C_ACCENT = RGBColor(91, 110, 135)
    C_TEXT = RGBColor(35, 39, 43)
    C_LINE = RGBColor(204, 212, 222)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_footer(slide, idx: int, total: int) -> None:
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.12), Inches(6.92), Inches(1.12), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(245, 248, 252)
        badge.line.fill.background()
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{idx}/{total}"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Times New Roman"
        p.font.size = Pt(11)
        p.font.bold = False
        p.font.color.rgb = C_ACCENT

    def add_header(slide, title: str) -> None:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = C_BG

        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.58))
        top.fill.solid()
        top.fill.fore_color.rgb = C_NAVY
        top.line.fill.background()

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.11), Inches(12.02), Inches(0.014))
        line.fill.solid()
        line.fill.fore_color.rgb = C_LINE
        line.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.73), Inches(11.6), Inches(0.35))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(25)
        p.font.bold = True
        p.font.color.rgb = C_NAVY

    def add_bullets(slide, bullets: List[str], x=0.95, y=1.38, w=11.3, h=5.7, size=22) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(size)
            p.font.color.rgb = C_TEXT
            p.space_after = Pt(10)
            p.line_spacing = 1.25

    def add_summary_band(slide, text: str, x: float, y: float, w: float, h: float) -> None:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(248, 251, 254)
        band.line.color.rgb = C_LINE
        band.line.width = Pt(0.8)
        tx = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.08), Inches(w - 0.24), Inches(h - 0.14))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(14.5)
        p.font.color.rgb = C_TEXT

    def add_vas_chart(slide, x: float, y: float, w: float, h: float) -> None:
        if not chart_data or "vas" not in chart_data:
            return
        d = chart_data["vas"]
        chart_d = CategoryChartData()
        chart_d.categories = d["categories"]
        chart_d.add_series("对照组", d["control"])
        chart_d.add_series("观察组", d["observation"])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), chart_d
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "VAS变化趋势（分）"
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = 6
        chart.value_axis.major_unit = 1
        chart.value_axis.has_major_gridlines = True

    def add_satisfaction_chart(slide, x: float, y: float, w: float, h: float) -> None:
        if not chart_data or "satisfaction" not in chart_data:
            return
        d = chart_data["satisfaction"]
        chart_d = CategoryChartData()
        chart_d.categories = d["categories"]
        chart_d.add_series("对照组", d["control"])
        chart_d.add_series("观察组", d["observation"])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_d
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "护理总满意度（%）"
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = 100
        chart.value_axis.major_unit = 20
        chart.value_axis.has_major_gridlines = True

    def add_compliance_chart(slide, x: float, y: float, w: float, h: float) -> None:
        if not chart_data or "compliance" not in chart_data:
            return
        d = chart_data["compliance"]
        chart_d = CategoryChartData()
        chart_d.categories = d["categories"]
        chart_d.add_series("对照组", d["control"])
        chart_d.add_series("观察组", d["observation"])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_d
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "依从性分项对比（分）"
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = 10
        chart.value_axis.major_unit = 2
        chart.value_axis.has_major_gridlines = True

    total = len(slides)
    for i, item in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        title = str(item["title"])
        bullets = [str(x) for x in item.get("bullets", [])]

        if i == 1:
            # Cover slide: restrained academic title page.
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = C_BG

            top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.76))
            top.fill.solid()
            top.fill.fore_color.rgb = C_NAVY
            top.line.fill.background()

            left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(0.09), Inches(3.05))
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = C_ACCENT
            left_bar.line.fill.background()

            t = slide.shapes.add_textbox(Inches(1.05), Inches(1.08), Inches(11.5), Inches(3.35))
            tf = t.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = "Microsoft YaHei"
            p0.font.size = Pt(40)
            p0.font.bold = True
            p0.font.color.rgb = C_NAVY
            p0.line_spacing = 1.12

            p1 = tf.add_paragraph()
            p1.text = "硕士复试学术汇报"
            p1.font.name = "Times New Roman"
            p1.font.size = Pt(21)
            p1.font.color.rgb = C_ACCENT
            p1.space_before = Pt(20)

            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.8), Inches(0.014))
            line.fill.solid()
            line.fill.fore_color.rgb = C_LINE
            line.line.fill.background()

            info = slide.shapes.add_textbox(Inches(0.85), Inches(5.82), Inches(11.8), Inches(1.0))
            itf = info.text_frame
            itf.clear()
            pi = itf.paragraphs[0]
            pi.text = f"汇报人：{presenter}    |    {affiliation}"
            pi.font.name = "Microsoft YaHei"
            pi.font.size = Pt(18)
            pi.font.bold = False
            pi.font.color.rgb = C_TEXT
            add_footer(slide, i, total)
            continue

        add_header(slide, title)

        if title == "关键数值：疼痛与满意度":
            if chart_data and "vas" in chart_data and "satisfaction" in chart_data:
                vas_c = chart_data["vas"]["control"]
                vas_o = chart_data["vas"]["observation"]
                sat_c = chart_data["satisfaction"]["control"][0]
                sat_o = chart_data["satisfaction"]["observation"][0]
                summary = (
                    f"VAS趋势：对照组 {vas_c[0]:.2f}→{vas_c[1]:.2f}→{vas_c[2]:.2f}；"
                    f"观察组 {vas_o[0]:.2f}→{vas_o[1]:.2f}→{vas_o[2]:.2f}。"
                    f" 护理总满意度：对照组 {sat_c:.0f}% vs 观察组 {sat_o:.0f}%。"
                )
            else:
                summary = "关键结果：观察组在疼痛缓解与护理满意度方面均优于对照组。"
            add_summary_band(slide, summary, x=0.95, y=1.2, w=11.3, h=0.56)
            add_vas_chart(slide, x=0.95, y=1.95, w=7.35, h=4.6)
            add_satisfaction_chart(slide, x=8.5, y=2.12, w=3.75, h=4.3)
        elif title == "依从性与行为改变":
            summary = "依从性结果：观察组在情绪管理、饮食控制、戒烟戒酒、及时用药、功能锻炼五项得分均高于对照组（P<0.05）。"
            add_summary_band(slide, summary, x=0.95, y=1.24, w=11.3, h=0.58)
            add_compliance_chart(slide, x=0.95, y=2.32, w=11.3, h=3.9)
        else:
            add_bullets(slide, bullets)

        add_footer(slide, i, total)

    prs.save(str(out_path))


def write_json(path: Path, data: object) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build submission PPT assets from paper file.")
    parser.add_argument("--source", required=True, help="Paper source path (.docx/.pdf/.txt/.md)")
    parser.add_argument("--output-dir", required=True, help="Output directory")
    parser.add_argument("--presenter", default="汪妍", help="Presenter name")
    parser.add_argument("--affiliation", default="南京大学医学院附属鼓楼医院 泌尿外科", help="Affiliation")
    args = parser.parse_args()

    source = Path(args.source).expanduser().resolve()
    out_dir = Path(args.output_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    markdown_text = extract_markdown(source)
    title = first_title(markdown_text)
    abstract_parts = extract_abstract_parts(markdown_text)
    numbers = extract_core_numbers(markdown_text)
    chart_data = extract_chart_data(numbers)
    slides = build_slides(title, abstract_parts, numbers, args.presenter, args.affiliation)
    requests = build_workspace_requests(slides)

    # 1) extracted source
    (out_dir / "paper_extracted.md").write_text(markdown_text, encoding="utf-8")

    # 2) outline and deck json
    outline = render_outline_markdown(title, args.presenter, args.affiliation, abstract_parts, slides)
    (out_dir / "submission_outline.md").write_text(outline, encoding="utf-8")
    write_json(out_dir / "submission_slides.json", {"title": title, "slides": slides})

    # 3) workspace-mcp payload templates
    write_json(
        out_dir / "workspace_mcp_create_args.json",
        {"user_google_email": "YOUR_GOOGLE_EMAIL", "title": f"{title}（复试汇报）"},
    )
    write_json(
        out_dir / "workspace_mcp_batch_update_args.template.json",
        {
            "user_google_email": "YOUR_GOOGLE_EMAIL",
            "presentation_id": "REPLACE_WITH_PRESENTATION_ID",
            "requests": requests,
        },
    )

    # 4) local draft pptx from markdown (fallback)
    slide_md = render_pandoc_slide_markdown(slides)
    slide_md_path = out_dir / "submission_slides.md"
    slide_md_path.write_text(slide_md, encoding="utf-8")
    subprocess.run(
        ["pandoc", str(slide_md_path), "-t", "pptx", "-o", str(out_dir / "submission_local_draft.pptx")],
        check=True,
    )
    # Preferred deck with designed theme.
    try:
        render_styled_pptx(
            slides,
            out_dir / "submission_local_styled.pptx",
            presenter=args.presenter,
            affiliation=args.affiliation,
            chart_data=chart_data,
        )
        # Alias: academic naming for direct use.
        render_styled_pptx(
            slides,
            out_dir / "submission_local_academic.pptx",
            presenter=args.presenter,
            affiliation=args.affiliation,
            chart_data=chart_data,
        )
    except ModuleNotFoundError:
        (out_dir / "submission_local_styled_build_note.txt").write_text(
            "python-pptx is missing. Run this script with:\n"
            "uv run --with python-pptx python3.10 build_submission_ppt_assets.py ...\n",
            encoding="utf-8",
        )

    # 5) helper script for workspace-mcp cli
    helper = textwrap.dedent(
        """\
        #!/usr/bin/env bash
        set -euo pipefail

        if [ "$#" -lt 1 ]; then
          echo "Usage:"
          echo "  $0 create <your_google_email>"
          echo "  $0 apply  <your_google_email> <presentation_id>"
          exit 1
        fi

        mode="$1"
        shift
        base_dir="$(cd "$(dirname "$0")" && pwd)"

        if [ "$mode" = "create" ]; then
          email="${1:?missing email}"
          python3.10 - "$email" "$base_dir" <<'PY'
        import json, sys, pathlib
        email = sys.argv[1]
        base = pathlib.Path(sys.argv[2])
        p = base / "workspace_mcp_create_args.json"
        data = json.loads(p.read_text(encoding="utf-8"))
        data["user_google_email"] = email
        print(json.dumps(data, ensure_ascii=False))
        PY
          exit 0
        fi

        if [ "$mode" = "apply" ]; then
          email="${1:?missing email}"
          pres_id="${2:?missing presentation_id}"
          python3.10 - "$email" "$pres_id" "$base_dir" <<'PY'
        import json, sys, pathlib
        email = sys.argv[1]
        pres = sys.argv[2]
        base = pathlib.Path(sys.argv[3])
        p = base / "workspace_mcp_batch_update_args.template.json"
        data = json.loads(p.read_text(encoding="utf-8"))
        data["user_google_email"] = email
        data["presentation_id"] = pres
        print(json.dumps(data, ensure_ascii=False))
        PY
          exit 0
        fi

        echo "Unknown mode: $mode"
        exit 1
        """
    )
    helper_path = out_dir / "workspace_mcp_payload.sh"
    helper_path.write_text(helper, encoding="utf-8")
    helper_path.chmod(0o755)

    readme = textwrap.dedent(
        """\
        # 投稿汇报PPT自动化产物

        本目录由 `build_submission_ppt_assets.py` 自动生成，包含：
        - `submission_outline.md`：论文汇报提纲（问题-方法-结果-结论-Q&A）
        - `submission_slides.json`：结构化幻灯片数据（12页）
        - `submission_slides.md`：Pandoc 幻灯片 Markdown
        - `submission_local_academic.pptx`：本地学术风格PPT（优先使用）
        - `submission_local_styled.pptx`：本地学术风格PPT（同内容别名）
        - `submission_local_draft.pptx`：本地基础版PPT（fallback）
        - `workspace_mcp_create_args.json`：`create_presentation` 参数模板
        - `workspace_mcp_batch_update_args.template.json`：`batch_update_presentation` 参数模板
        - `workspace_mcp_env.example`：Google OAuth 环境变量模板
        - `workspace_mcp_payload.sh`：根据邮箱和 presentation_id 生成可直接 pipe 给 CLI 的 JSON

        ## 用 Google Workspace MCP 创建正式 Google Slides

        0. 先配置 Google OAuth 环境变量（必需）：

        ```bash
        export GOOGLE_OAUTH_CLIENT_ID="YOUR_CLIENT_ID"
        export GOOGLE_OAUTH_CLIENT_SECRET="YOUR_CLIENT_SECRET"
        ```

        1. 先创建演示文稿（得到 presentation_id）：

        ```bash
        ./workspace_mcp_payload.sh create your_email@example.com \\
          | uvx workspace-mcp --cli create_presentation
        ```

        2. 把第1步输出中的 `presentation_id` 复制出来后，批量写入全部页面：

        ```bash
        ./workspace_mcp_payload.sh apply your_email@example.com PRESENTATION_ID \\
          | uvx workspace-mcp --cli batch_update_presentation
        ```

        3. 如需二次修改页面内容，编辑 `submission_slides.json` 后重新运行生成脚本。
        """
    )
    (out_dir / "README.md").write_text(readme, encoding="utf-8")
    (out_dir / "workspace_mcp_env.example").write_text(
        textwrap.dedent(
            """\
            # Copy and fill these env vars before calling workspace-mcp CLI.
            export GOOGLE_OAUTH_CLIENT_ID="YOUR_CLIENT_ID"
            export GOOGLE_OAUTH_CLIENT_SECRET="YOUR_CLIENT_SECRET"
            """
        ),
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
